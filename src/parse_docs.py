# src/parse_docs.py
import ssl
ssl._create_default_https_context = ssl._create_unverified_context
import os
import glob
import fitz  # PyMuPDF
from pptx import Presentation
from unstructured.partition.docx import partition_docx
import easyocr

# 初始化 OCR（中文+英文）
print("正在加载 EasyOCR（首次运行会下载模型）...")
reader = easyocr.Reader(['ch_sim', 'en'], verbose=False)

def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    full_text = ""
    for page_num, page in enumerate(doc):
        # 1. 原生文字
        text = page.get_text().strip()
        if text:
            full_text += f"\n--- Page {page_num + 1} (Text) ---\n{text}\n"
        
        # 2. 图片 OCR
        image_list = page.get_images(full=True)
        for img_idx, img in enumerate(image_list):
            try:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                ocr_results = reader.readtext(image_bytes, detail=0)
                ocr_text = "\n".join(ocr_results).strip()
                if ocr_text:
                    full_text += f"\n--- Page {page_num + 1} (Image {img_idx + 1} OCR) ---\n{ocr_text}\n"
            except Exception as e:
                print(f"  ⚠️ OCR 失败 (Page {page_num + 1}, Image {img_idx + 1})")
    doc.close()
    return full_text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    full_text = ""
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        if slide_text.strip():
            full_text += f"\n--- Slide {i + 1} ---\n{slide_text}\n"
    return full_text

def extract_text_from_docx(docx_path):
    elements = partition_docx(docx_path)
    return "\n\n".join([str(e) for e in elements])

def main():
    os.makedirs("processed_texts", exist_ok=True)
    
    for file_path in glob.glob("raw_docs/*"):
        filename = os.path.basename(file_path)
        print(f"📄 处理: {filename}")
        
        try:
            if filename.lower().endswith(".pdf"):
                text = extract_text_from_pdf(file_path)
            elif filename.lower().endswith(".pptx"):
                text = extract_text_from_pptx(file_path)
            elif filename.lower().endswith(".docx"):
                text = extract_text_from_docx(file_path)
            else:
                print(f"  ⚠️ 跳过不支持格式: {filename}")
                continue
            
            out_name = os.path.splitext(filename)[0] + ".txt"
            with open(f"processed_texts/{out_name}", "w", encoding="utf-8") as f:
                f.write(text)
            print(f"  ✅ 已保存: {out_name}")
        except Exception as e:
            print(f"  ❌ 处理失败: {e}")

if __name__ == "__main__":
    main()